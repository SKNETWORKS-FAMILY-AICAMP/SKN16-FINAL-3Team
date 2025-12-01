"""
하경은행 화면 설계서 PPT 템플릿 생성기
- 테마 색상: Primary(갈색), Bank(블루), Secondary(베이지)
- 사진 플레이스홀더 포함
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# 하경은행 테마 색상
COLORS = {
    'primary': RGBColor(0xD4, 0xA5, 0x74),      # #d4a574 - 메인 갈색
    'primary_dark': RGBColor(0x9C, 0x7F, 0x5F), # #9c7f5f - 진한 갈색
    'bank': RGBColor(0x62, 0x7D, 0x98),         # #627d98 - 은행 블루
    'bank_dark': RGBColor(0x33, 0x4E, 0x68),    # #334e68 - 진한 블루
    'secondary': RGBColor(0xCD, 0xC3, 0xB4),    # #cdc3b4 - 베이지
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x24, 0x3B, 0x53),        # #243b53 - 다크 블루
    'gray': RGBColor(0x82, 0x9A, 0xB1),         # #829ab1 - 그레이 블루
    'light_bg': RGBColor(0xF0, 0xF4, 0xF8),     # #f0f4f8 - 밝은 배경
}

# 화면 목록 데이터
SCREENS = [
    # AUTH - 인증
    {"id": "SCR-AUTH-01", "name": "로그인 페이지", "category": "AUTH", "func_id": "AUTH01", "desc": "이메일/사번 입력 폼, 로그인 버튼"},
    {"id": "SCR-AUTH-02", "name": "로그인 실패", "category": "AUTH", "func_id": "AUTH01-03,04", "desc": "에러 메시지 표시 상태"},
    {"id": "SCR-AUTH-03", "name": "아이디 찾기", "category": "AUTH", "func_id": "AUTH03-02", "desc": "이름, 사번 입력 폼 + 결과 화면"},
    {"id": "SCR-AUTH-04", "name": "비밀번호 재설정", "category": "AUTH", "func_id": "AUTH03-01", "desc": "이메일, 사번 입력 폼"},
    {"id": "SCR-AUTH-05", "name": "내 정보 조회", "category": "AUTH", "func_id": "AUTH02-01", "desc": "프로필 페이지 (이름, 부서, 사진)"},
    {"id": "SCR-AUTH-06", "name": "내 정보 수정", "category": "AUTH", "func_id": "AUTH02-02,04", "desc": "비밀번호 변경, 사진 업로드"},
    
    # TC - 연수원 관리
    {"id": "SCR-TC-01", "name": "연수원 DB 동기화", "category": "TC", "func_id": "TC01", "desc": "기수 선택 UI, 동기화 버튼"},
    {"id": "SCR-TC-02", "name": "동기화 결과", "category": "TC", "func_id": "TC01-03", "desc": "생성된 멘티/멘토 수 표시"},
    {"id": "SCR-TC-03", "name": "멘티 목록", "category": "TC", "func_id": "TC04-01,02", "desc": "목록 테이블, 검색창, 필터"},
    {"id": "SCR-TC-04", "name": "멘토 목록", "category": "TC", "func_id": "TC04-01", "desc": "멘토 목록 테이블"},
    {"id": "SCR-TC-05", "name": "레코드 삭제 확인", "category": "TC", "func_id": "TC04-03", "desc": "삭제 확인 모달"},
    
    # CHAT - AI 챗봇
    {"id": "SCR-CHAT-01", "name": "챗봇 메인", "category": "CHAT", "func_id": "CHAT01", "desc": "채팅 인터페이스 전체"},
    {"id": "SCR-CHAT-02", "name": "규정 질의응답", "category": "CHAT", "func_id": "CHAT01-02", "desc": "답변 + 출처(sources) 표시"},
    {"id": "SCR-CHAT-03", "name": "일정 생성 대화", "category": "CHAT", "func_id": "CHAT02", "desc": "자연어로 일정 생성"},
    {"id": "SCR-CHAT-04", "name": "학습 현황 분석", "category": "CHAT", "func_id": "CHAT03", "desc": "약점 분석 결과 답변"},
    
    # SCH - 일정 관리
    {"id": "SCR-SCH-01", "name": "캘린더 뷰 (월간)", "category": "SCH", "func_id": "SCH02-01", "desc": "월간 캘린더 + 일정 표시"},
    {"id": "SCR-SCH-02", "name": "캘린더 뷰 (주간)", "category": "SCH", "func_id": "SCH02-01", "desc": "주간 또는 일간 뷰"},
    {"id": "SCR-SCH-03", "name": "일정 등록 모달", "category": "SCH", "func_id": "SCH01-01", "desc": "제목, 시간, 장소, 색상 입력"},
    {"id": "SCR-SCH-04", "name": "일정 상세/수정", "category": "SCH", "func_id": "SCH02-03", "desc": "상세 정보 + 수정/삭제 버튼"},
    {"id": "SCR-SCH-05", "name": "공통 빈 시간 추천", "category": "SCH", "func_id": "SCH04-01", "desc": "멘토링 가능 시간 슬롯"},
    
    # SIM - RAG 시뮬레이션
    {"id": "SCR-SIM-01", "name": "시뮬레이션 선택", "category": "SIM", "func_id": "SIM01", "desc": "페르소나 + 상황 선택 UI"},
    {"id": "SCR-SIM-02", "name": "페르소나 카드", "category": "SIM", "func_id": "SIM01-01", "desc": "고객 유형 카드"},
    {"id": "SCR-SIM-03", "name": "상황 카드", "category": "SIM", "func_id": "SIM01-02", "desc": "시나리오 카드"},
    {"id": "SCR-SIM-04", "name": "대화 진행 화면", "category": "SIM", "func_id": "SIM03", "desc": "음성 대화 UI, 녹음 버튼"},
    {"id": "SCR-SIM-05", "name": "실시간 점수", "category": "SIM", "func_id": "SIM03-03", "desc": "대화 중 점수 표시"},
    {"id": "SCR-SIM-06", "name": "종합 피드백", "category": "SIM", "func_id": "SIM04", "desc": "6대 역량 점수 차트"},
    {"id": "SCR-SIM-07", "name": "피드백 이력", "category": "SIM", "func_id": "SIM04-03", "desc": "과거 피드백 목록"},
    
    # QUIZ - 퀴즈
    {"id": "SCR-QUIZ-01", "name": "퀴즈 모드 선택", "category": "QUIZ", "func_id": "QUIZ01", "desc": "랜덤/맞춤/평가 선택"},
    {"id": "SCR-QUIZ-02", "name": "퀴즈 풀이 화면", "category": "QUIZ", "func_id": "QUIZ01", "desc": "문제 + 보기 + 진행률"},
    {"id": "SCR-QUIZ-03", "name": "퀴즈 결과", "category": "QUIZ", "func_id": "QUIZ03-01", "desc": "점수, 정답률, 통계"},
    {"id": "SCR-QUIZ-04", "name": "오답 노트", "category": "QUIZ", "func_id": "QUIZ04-01", "desc": "틀린 문제 + 해설"},
    {"id": "SCR-QUIZ-05", "name": "응시 횟수 제한", "category": "QUIZ", "func_id": "QUIZ02-02", "desc": "횟수 초과 에러 메시지"},
    
    # REP - 대시보드/리포트
    {"id": "SCR-REP-01", "name": "멘티 대시보드", "category": "REP", "func_id": "REP01", "desc": "전체 레이아웃"},
    {"id": "SCR-REP-02", "name": "역량 레이더 차트", "category": "REP", "func_id": "REP01-03", "desc": "6대 역량 시각화"},
    {"id": "SCR-REP-03", "name": "성장 추이 그래프", "category": "REP", "func_id": "REP01-04", "desc": "시간별 점수 변화"},
    {"id": "SCR-REP-04", "name": "멘토 대시보드", "category": "REP", "func_id": "REP02", "desc": "담당 멘티 목록 + 현황"},
    {"id": "SCR-REP-05", "name": "관리자 매칭 현황", "category": "REP", "func_id": "REP03-02", "desc": "전체 매칭 테이블"},
    
    # CLUB - 동아리 라운지
    {"id": "SCR-CLUB-01", "name": "게시판 목록", "category": "CLUB", "func_id": "CLUB01-03", "desc": "게시글 리스트"},
    {"id": "SCR-CLUB-02", "name": "게시글 작성", "category": "CLUB", "func_id": "CLUB01-01", "desc": "제목, 내용, 카테고리 입력"},
    {"id": "SCR-CLUB-03", "name": "게시글 상세", "category": "CLUB", "func_id": "CLUB01-04", "desc": "본문 + 댓글 목록"},
    {"id": "SCR-CLUB-04", "name": "같이하기 승인", "category": "CLUB", "func_id": "CLUB02-03", "desc": "승인 버튼/완료 상태"},
]

CATEGORY_NAMES = {
    "AUTH": "인증",
    "TC": "연수원 관리",
    "CHAT": "AI 챗봇",
    "SCH": "일정 관리",
    "SIM": "RAG 시뮬레이션",
    "QUIZ": "퀴즈",
    "REP": "대시보드/리포트",
    "CLUB": "동아리 라운지",
}


def add_title_slide(prs):
    """표지 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경색 설정 (bank 색상)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bank_dark']
    background.line.fill.background()
    
    # 상단 장식 바
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['primary']
    top_bar.line.fill.background()
    
    # 메인 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "하경은행 CANT 시스템"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 서브 타이틀
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "화면 설계서"
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # 하단 정보
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Comprehensive AI-based New employee Training"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER


def add_toc_slide(prs):
    """목차 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['light_bg']
    background.line.fill.background()
    
    # 헤더 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['bank_dark']
    header.line.fill.background()
    
    # 타이틀
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 목차 (Table of Contents)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 목차 내용
    categories = list(CATEGORY_NAMES.items())
    
    for i, (cat_id, cat_name) in enumerate(categories):
        count = len([s for s in SCREENS if s['category'] == cat_id])
        
        row = i // 2
        col = i % 2
        
        x = Inches(0.8 + col * 4.5)
        y = Inches(1.6 + row * 0.7)
        
        # 번호 원
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['primary']
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(x, y + Inches(0.05), Inches(0.4), Inches(0.3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # 카테고리 이름
        text_box = slide.shapes.add_textbox(x + Inches(0.5), y, Inches(3.5), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{cat_name} ({count}개 화면)"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['black']


def add_section_slide(prs, category_id, category_name):
    """섹션 구분 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bank']
    background.line.fill.background()
    
    # 왼쪽 장식 바
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = COLORS['primary']
    left_bar.line.fill.background()
    
    # 카테고리 ID
    id_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    tf = id_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"[ {category_id} ]"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['primary']
    
    # 카테고리 이름
    name_box = slide.shapes.add_textbox(Inches(1), Inches(3.1), Inches(8), Inches(1))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = category_name
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 화면 수
    count = len([s for s in SCREENS if s['category'] == category_id])
    count_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
    tf = count_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{count}개 화면"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['secondary']


def add_screen_slide(prs, screen):
    """개별 화면 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 상단 헤더 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['bank_dark']
    header.line.fill.background()
    
    # 화면 ID
    id_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(2), Inches(0.35))
    tf = id_box.text_frame
    p = tf.paragraphs[0]
    p.text = screen['id']
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['primary']
    
    # 화면 이름
    name_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.45), Inches(6), Inches(0.5))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = screen['name']
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 기능 ID 태그
    func_tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(0.3), Inches(2.2), Inches(0.4)
    )
    func_tag.fill.solid()
    func_tag.fill.fore_color.rgb = COLORS['primary']
    func_tag.line.fill.background()
    
    func_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.35), Inches(2.2), Inches(0.35))
    tf = func_box.text_frame
    p = tf.paragraphs[0]
    p.text = screen['func_id']
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 스크린샷 플레이스홀더 영역
    screenshot_area = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.2), Inches(6.4), Inches(4.3)
    )
    screenshot_area.fill.solid()
    screenshot_area.fill.fore_color.rgb = COLORS['light_bg']
    screenshot_area.line.color.rgb = COLORS['gray']
    screenshot_area.line.width = Pt(1)
    
    # 플레이스홀더 텍스트
    placeholder_text = slide.shapes.add_textbox(Inches(0.3), Inches(3.0), Inches(6.4), Inches(0.8))
    tf = placeholder_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📷 스크린샷을 여기에 삽입하세요"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER
    
    # 우측 정보 패널
    info_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.2), Inches(2.9), Inches(4.3)
    )
    info_panel.fill.solid()
    info_panel.fill.fore_color.rgb = COLORS['light_bg']
    info_panel.line.fill.background()
    
    # 정보 패널 - 화면 설명
    desc_label = slide.shapes.add_textbox(Inches(7.1), Inches(1.4), Inches(2.5), Inches(0.3))
    tf = desc_label.text_frame
    p = tf.paragraphs[0]
    p.text = "📝 화면 설명"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['bank_dark']
    
    desc_box = slide.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(2.5), Inches(1.0))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = screen['desc']
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['black']
    
    # 정보 패널 - 연관 기능
    func_label = slide.shapes.add_textbox(Inches(7.1), Inches(2.9), Inches(2.5), Inches(0.3))
    tf = func_label.text_frame
    p = tf.paragraphs[0]
    p.text = "🔗 연관 기능 ID"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['bank_dark']
    
    func_value = slide.shapes.add_textbox(Inches(7.1), Inches(3.2), Inches(2.5), Inches(0.4))
    tf = func_value.text_frame
    p = tf.paragraphs[0]
    p.text = screen['func_id']
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['primary_dark']
    
    # 정보 패널 - 비고
    note_label = slide.shapes.add_textbox(Inches(7.1), Inches(3.8), Inches(2.5), Inches(0.3))
    tf = note_label.text_frame
    p = tf.paragraphs[0]
    p.text = "📌 비고"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLORS['bank_dark']
    
    note_box = slide.shapes.add_textbox(Inches(7.1), Inches(4.1), Inches(2.5), Inches(1.2))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "(여기에 추가 설명을 작성하세요)"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS['gray']


def create_ppt():
    """PPT 생성 메인 함수"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. 표지
    add_title_slide(prs)
    
    # 2. 목차
    add_toc_slide(prs)
    
    # 3. 각 카테고리별 섹션 + 화면 슬라이드
    current_category = None
    for screen in SCREENS:
        # 새로운 카테고리면 섹션 슬라이드 추가
        if screen['category'] != current_category:
            current_category = screen['category']
            add_section_slide(prs, current_category, CATEGORY_NAMES[current_category])
        
        # 화면 슬라이드 추가
        add_screen_slide(prs, screen)
    
    # 저장
    os.makedirs("docs/standardized", exist_ok=True)
    file_path = "docs/standardized/CANT_화면설계서_Template.pptx"
    prs.save(file_path)
    
    print(f"✅ PPT 템플릿 생성 완료: {file_path}")
    print(f"   - 총 슬라이드 수: {len(prs.slides)}장")
    print(f"   - 표지: 1장")
    print(f"   - 목차: 1장")
    print(f"   - 섹션 구분: {len(CATEGORY_NAMES)}장")
    print(f"   - 화면 슬라이드: {len(SCREENS)}장")


if __name__ == "__main__":
    create_ppt()

